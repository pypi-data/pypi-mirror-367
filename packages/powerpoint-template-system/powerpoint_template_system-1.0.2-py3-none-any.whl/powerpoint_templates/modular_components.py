"""
Modular Component System for PowerPoint Generation

This module provides reusable, configurable components that can be mixed and matched
to create different presentation styles and layouts.
"""

from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Any, Tuple, Callable
from enum import Enum
import json
from pathlib import Path
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class ComponentType(Enum):
    """Types of presentation components"""
    HEADER = "header"
    FOOTER = "footer"
    TITLE = "title"
    CONTENT = "content"
    SIDEBAR = "sidebar"
    CHART = "chart"
    IMAGE = "image"
    TABLE = "table"
    BULLET_LIST = "bullet_list"
    QUOTE = "quote"
    LOGO = "logo"
    DIVIDER = "divider"


class ComponentPosition(Enum):
    """Component positioning options"""
    TOP = "top"
    BOTTOM = "bottom"
    LEFT = "left"
    RIGHT = "right"
    CENTER = "center"
    TOP_LEFT = "top_left"
    TOP_RIGHT = "top_right"
    BOTTOM_LEFT = "bottom_left"
    BOTTOM_RIGHT = "bottom_right"
    FULL = "full"


@dataclass
class ComponentBounds:
    """Defines the bounds of a component on a slide"""
    left: float  # inches from left edge
    top: float   # inches from top edge
    width: float # width in inches
    height: float # height in inches
    
    def to_inches(self) -> Tuple[int, int, int, int]:
        """Convert to python-pptx Inches format"""
        return (Inches(self.left), Inches(self.top), 
                Inches(self.width), Inches(self.height))


@dataclass
class ComponentStyle:
    """Styling configuration for components"""
    font_name: str = "Calibri"
    font_size: int = 12
    font_color: str = "#000000"
    font_bold: bool = False
    font_italic: bool = False
    background_color: Optional[str] = None
    border_color: Optional[str] = None
    border_width: float = 0
    alignment: str = "left"  # left, center, right, justify
    padding: float = 0.1  # inches
    margin: float = 0.05  # inches
    
    def get_rgb_color(self, color_hex: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        color_hex = color_hex.lstrip('#')
        return RGBColor(
            int(color_hex[0:2], 16),
            int(color_hex[2:4], 16),
            int(color_hex[4:6], 16)
        )
    
    def get_alignment(self) -> PP_ALIGN:
        """Get python-pptx alignment constant"""
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }
        return alignment_map.get(self.alignment, PP_ALIGN.LEFT)


class PresentationComponent(ABC):
    """Abstract base class for all presentation components"""
    
    def __init__(self, component_id: str, component_type: ComponentType, 
                 bounds: ComponentBounds, style: ComponentStyle):
        self.component_id = component_id
        self.component_type = component_type
        self.bounds = bounds
        self.style = style
        self.visible = True
        self.z_order = 0  # Higher values appear on top
    
    @abstractmethod
    def render(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render the component on the slide"""
        pass
    
    @abstractmethod
    def validate_data(self, data: Dict[str, Any]) -> bool:
        """Validate the data required for this component"""
        pass
    
    def set_visibility(self, visible: bool) -> None:
        """Set component visibility"""
        self.visible = visible
    
    def set_z_order(self, z_order: int) -> None:
        """Set component z-order (layering)"""
        self.z_order = z_order


class HeaderComponent(PresentationComponent):
    """Reusable header component"""
    
    def __init__(self, component_id: str = "header", 
                 bounds: ComponentBounds = None, 
                 style: ComponentStyle = None):
        if bounds is None:
            bounds = ComponentBounds(0, 0, 10, 1)
        if style is None:
            style = ComponentStyle(font_size=16, font_bold=True)
        
        super().__init__(component_id, ComponentType.HEADER, bounds, style)
        self.show_logo = True
        self.show_title = True
        self.show_subtitle = False
        self.show_date = False
        self.logo_size = 0.8  # inches
    
    def render(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render header component"""
        if not self.visible:
            return None
        
        left, top, width, height = self.bounds.to_inches()
        
        # Add background if specified
        if self.style.background_color:
            bg_shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.style.get_rgb_color(self.style.background_color)
            bg_shape.line.fill.background()
        
        # Add logo if specified
        logo_width = 0
        if self.show_logo and data.get('logo_path'):
            try:
                logo = slide.shapes.add_picture(
                    data['logo_path'],
                    left + Inches(self.style.padding),
                    top + Inches(self.style.padding),
                    height=Inches(self.logo_size)
                )
                logo_width = logo.width
            except Exception as e:
                print(f"Warning: Could not add logo: {e}")
        
        # Add title text
        if self.show_title and data.get('title'):
            title_left = left + logo_width + Inches(self.style.padding * 2)
            title_width = width - logo_width - Inches(self.style.padding * 3)
            
            title_box = slide.shapes.add_textbox(
                title_left, top + Inches(self.style.padding),
                title_width, height - Inches(self.style.padding * 2)
            )
            
            title_frame = title_box.text_frame
            title_frame.text = data['title']
            title_frame.paragraphs[0].font.name = self.style.font_name
            title_frame.paragraphs[0].font.size = Pt(self.style.font_size)
            title_frame.paragraphs[0].font.bold = self.style.font_bold
            title_frame.paragraphs[0].font.color.rgb = self.style.get_rgb_color(self.style.font_color)
            title_frame.paragraphs[0].alignment = self.style.get_alignment()
        
        return True
    
    def validate_data(self, data: Dict[str, Any]) -> bool:
        """Validate header data"""
        if self.show_title and 'title' not in data:
            return False
        if self.show_logo and 'logo_path' not in data:
            return False
        return True


class FooterComponent(PresentationComponent):
    """Reusable footer component"""
    
    def __init__(self, component_id: str = "footer",
                 bounds: ComponentBounds = None,
                 style: ComponentStyle = None):
        if bounds is None:
            bounds = ComponentBounds(0, 6.5, 10, 0.5)
        if style is None:
            style = ComponentStyle(font_size=10, alignment="center")
        
        super().__init__(component_id, ComponentType.FOOTER, bounds, style)
        self.show_page_number = True
        self.show_company_name = True
        self.show_confidentiality = True
        self.confidentiality_text = "Confidential"
    
    def render(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render footer component"""
        if not self.visible:
            return None
        
        left, top, width, height = self.bounds.to_inches()
        
        # Add background if specified
        if self.style.background_color:
            bg_shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.style.get_rgb_color(self.style.background_color)
            bg_shape.line.fill.background()
        
        # Build footer text
        footer_parts = []
        
        if self.show_company_name and data.get('company_name'):
            footer_parts.append(data['company_name'])
        
        if self.show_confidentiality:
            footer_parts.append(self.confidentiality_text)
        
        if self.show_page_number and data.get('page_number'):
            footer_parts.append(f"Page {data['page_number']}")
        
        if footer_parts:
            footer_text = " | ".join(footer_parts)
            
            footer_box = slide.shapes.add_textbox(
                left + Inches(self.style.padding),
                top + Inches(self.style.padding),
                width - Inches(self.style.padding * 2),
                height - Inches(self.style.padding * 2)
            )
            
            footer_frame = footer_box.text_frame
            footer_frame.text = footer_text
            footer_frame.paragraphs[0].font.name = self.style.font_name
            footer_frame.paragraphs[0].font.size = Pt(self.style.font_size)
            footer_frame.paragraphs[0].font.color.rgb = self.style.get_rgb_color(self.style.font_color)
            footer_frame.paragraphs[0].alignment = self.style.get_alignment()
        
        return True
    
    def validate_data(self, data: Dict[str, Any]) -> bool:
        """Validate footer data"""
        return True  # Footer is generally self-contained


class ContentComponent(PresentationComponent):
    """Flexible content component for various content types"""
    
    def __init__(self, component_id: str = "content",
                 bounds: ComponentBounds = None,
                 style: ComponentStyle = None):
        if bounds is None:
            bounds = ComponentBounds(0.5, 1.5, 9, 4.5)
        if style is None:
            style = ComponentStyle()
        
        super().__init__(component_id, ComponentType.CONTENT, bounds, style)
        self.content_type = "text"  # text, bullet_list, chart, image, table
        self.max_bullet_levels = 3
        self.bullet_indent = 0.25  # inches per level
    
    def render(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render content based on content type"""
        if not self.visible:
            return None
        
        content_type = data.get('content_type', self.content_type)
        
        render_methods = {
            'text': self._render_text,
            'bullet_list': self._render_bullet_list,
            'chart': self._render_chart_placeholder,
            'image': self._render_image,
            'table': self._render_table_placeholder
        }
        
        render_method = render_methods.get(content_type, self._render_text)
        return render_method(slide, data)
    
    def _render_text(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render plain text content"""
        left, top, width, height = self.bounds.to_inches()
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = data.get('text', '')
        
        # Apply styling
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.style.font_name
        paragraph.font.size = Pt(self.style.font_size)
        paragraph.font.bold = self.style.font_bold
        paragraph.font.italic = self.style.font_italic
        paragraph.font.color.rgb = self.style.get_rgb_color(self.style.font_color)
        paragraph.alignment = self.style.get_alignment()
        
        return text_box
    
    def _render_bullet_list(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render bullet list content"""
        left, top, width, height = self.bounds.to_inches()
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        items = data.get('items', [])
        for i, item in enumerate(items):
            if isinstance(item, dict):
                text = item.get('text', '')
                level = min(item.get('level', 0), self.max_bullet_levels - 1)
            else:
                text = str(item)
                level = 0
            
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = text
            paragraph.level = level
            paragraph.font.name = self.style.font_name
            paragraph.font.size = Pt(max(8, self.style.font_size - level * 2))
            paragraph.font.color.rgb = self.style.get_rgb_color(self.style.font_color)
        
        return text_box
    
    def _render_chart_placeholder(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render chart placeholder (to be implemented with actual charting)"""
        left, top, width, height = self.bounds.to_inches()
        
        # For now, create a placeholder rectangle
        chart_placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        chart_placeholder.fill.solid()
        chart_placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        chart_placeholder.line.color.rgb = RGBColor(128, 128, 128)
        
        # Add chart title
        title_box = slide.shapes.add_textbox(
            left, top + Inches(0.1), width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = data.get('chart_title', 'Chart Placeholder')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.bold = True
        
        return chart_placeholder
    
    def _render_image(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render image content"""
        image_path = data.get('image_path')
        if not image_path:
            return None
        
        try:
            left, top, width, height = self.bounds.to_inches()
            image = slide.shapes.add_picture(image_path, left, top, width, height)
            return image
        except Exception as e:
            print(f"Warning: Could not add image: {e}")
            return None
    
    def _render_table_placeholder(self, slide: Slide, data: Dict[str, Any]) -> Any:
        """Render table placeholder (to be implemented with actual table creation)"""
        left, top, width, height = self.bounds.to_inches()
        
        # Create placeholder for table
        table_placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        table_placeholder.fill.solid()
        table_placeholder.fill.fore_color.rgb = RGBColor(250, 250, 250)
        table_placeholder.line.color.rgb = RGBColor(128, 128, 128)
        
        # Add table title
        title_box = slide.shapes.add_textbox(
            left, top + Inches(0.1), width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = data.get('table_title', 'Table Placeholder')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.bold = True
        
        return table_placeholder
    
    def validate_data(self, data: Dict[str, Any]) -> bool:
        """Validate content data"""
        content_type = data.get('content_type', self.content_type)
        
        if content_type == 'text' and 'text' not in data:
            return False
        elif content_type == 'bullet_list' and 'items' not in data:
            return False
        elif content_type == 'image' and 'image_path' not in data:
            return False
        
        return True


class ComponentLayout:
    """Manages the layout and positioning of components on a slide"""
    
    def __init__(self, slide_width: float = 10, slide_height: float = 7.5):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.components: List[PresentationComponent] = []
        self.layout_name = "custom"
    
    def add_component(self, component: PresentationComponent) -> None:
        """Add a component to the layout"""
        self.components.append(component)
        # Sort by z-order to ensure proper layering
        self.components.sort(key=lambda c: c.z_order)
    
    def remove_component(self, component_id: str) -> bool:
        """Remove a component by ID"""
        for i, component in enumerate(self.components):
            if component.component_id == component_id:
                del self.components[i]
                return True
        return False
    
    def get_component(self, component_id: str) -> Optional[PresentationComponent]:
        """Get a component by ID"""
        for component in self.components:
            if component.component_id == component_id:
                return component
        return None
    
    def render_layout(self, slide: Slide, data: Dict[str, Any]) -> bool:
        """Render all components in the layout"""
        success = True
        
        for component in self.components:
            try:
                component_data = data.get(component.component_id, {})
                if component.validate_data(component_data):
                    component.render(slide, component_data)
                else:
                    print(f"Warning: Invalid data for component {component.component_id}")
                    success = False
            except Exception as e:
                print(f"Error rendering component {component.component_id}: {e}")
                success = False
        
        return success
    
    def auto_layout_two_column(self) -> None:
        """Automatically arrange components in a two-column layout"""
        content_components = [c for c in self.components 
                            if c.component_type == ComponentType.CONTENT]
        
        if len(content_components) >= 2:
            # Left column
            content_components[0].bounds = ComponentBounds(0.5, 1.5, 4.5, 4.5)
            # Right column
            content_components[1].bounds = ComponentBounds(5.25, 1.5, 4.25, 4.5)
    
    def auto_layout_three_column(self) -> None:
        """Automatically arrange components in a three-column layout"""
        content_components = [c for c in self.components 
                            if c.component_type == ComponentType.CONTENT]
        
        if len(content_components) >= 3:
            column_width = 3.0
            for i, component in enumerate(content_components[:3]):
                component.bounds = ComponentBounds(
                    0.5 + i * (column_width + 0.25), 1.5, column_width, 4.5
                )
    
    def export_layout(self, file_path: str) -> None:
        """Export layout configuration to JSON"""
        layout_data = {
            'layout_name': self.layout_name,
            'slide_width': self.slide_width,
            'slide_height': self.slide_height,
            'components': []
        }
        
        for component in self.components:
            component_data = {
                'component_id': component.component_id,
                'component_type': component.component_type.value,
                'bounds': {
                    'left': component.bounds.left,
                    'top': component.bounds.top,
                    'width': component.bounds.width,
                    'height': component.bounds.height
                },
                'style': component.style.__dict__,
                'visible': component.visible,
                'z_order': component.z_order
            }
            layout_data['components'].append(component_data)
        
        with open(file_path, 'w') as f:
            json.dump(layout_data, f, indent=2)


class LayoutLibrary:
    """Library of predefined component layouts"""
    
    @staticmethod
    def create_standard_business_layout() -> ComponentLayout:
        """Create a standard business presentation layout"""
        layout = ComponentLayout()
        layout.layout_name = "standard_business"
        
        # Header
        header = HeaderComponent(
            bounds=ComponentBounds(0, 0, 10, 1),
            style=ComponentStyle(font_size=18, font_bold=True, alignment="left")
        )
        layout.add_component(header)
        
        # Content
        content = ContentComponent(
            bounds=ComponentBounds(0.5, 1.5, 9, 4.5),
            style=ComponentStyle(font_size=14)
        )
        layout.add_component(content)
        
        # Footer
        footer = FooterComponent(
            bounds=ComponentBounds(0, 6.5, 10, 0.5),
            style=ComponentStyle(font_size=10, alignment="center")
        )
        layout.add_component(footer)
        
        return layout
    
    @staticmethod
    def create_executive_layout() -> ComponentLayout:
        """Create an executive presentation layout"""
        layout = ComponentLayout()
        layout.layout_name = "executive"
        
        # Header with larger logo space
        header = HeaderComponent(
            bounds=ComponentBounds(0, 0, 10, 1.2),
            style=ComponentStyle(font_size=20, font_bold=True, alignment="left")
        )
        header.logo_size = 1.0
        layout.add_component(header)
        
        # Main content
        content = ContentComponent(
            bounds=ComponentBounds(0.5, 1.5, 9, 4.5),
            style=ComponentStyle(font_size=16)
        )
        layout.add_component(content)
        
        # Footer with confidentiality emphasis
        footer = FooterComponent(
            bounds=ComponentBounds(0, 6.5, 10, 0.5),
            style=ComponentStyle(font_size=10, alignment="center", font_bold=True)
        )
        footer.confidentiality_text = "CONFIDENTIAL - Executive Summary"
        layout.add_component(footer)
        
        return layout
    
    @staticmethod
    def create_two_column_layout() -> ComponentLayout:
        """Create a two-column layout"""
        layout = ComponentLayout()
        layout.layout_name = "two_column"
        
        # Header
        header = HeaderComponent()
        layout.add_component(header)
        
        # Left content
        left_content = ContentComponent(
            component_id="left_content",
            bounds=ComponentBounds(0.5, 1.5, 4.5, 4.5)
        )
        layout.add_component(left_content)
        
        # Right content
        right_content = ContentComponent(
            component_id="right_content",
            bounds=ComponentBounds(5.25, 1.5, 4.25, 4.5)
        )
        layout.add_component(right_content)
        
        # Footer
        footer = FooterComponent()
        layout.add_component(footer)
        
        return layout


# Usage example
if __name__ == "__main__":
    # Create a presentation with modular components
    from pptx import Presentation
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Create layout
    layout = LayoutLibrary.create_standard_business_layout()
    
    # Sample data
    slide_data = {
        'header': {
            'title': 'Business Overview',
            'logo_path': 'assets/logo.png'
        },
        'content': {
            'content_type': 'bullet_list',
            'items': [
                'Market Analysis',
                'Financial Performance',
                {'text': 'Strategic Initiatives', 'level': 1},
                {'text': 'Q1 Goals', 'level': 1}
            ]
        },
        'footer': {
            'company_name': 'Acme Corporation',
            'page_number': 1
        }
    }
    
    # Render the layout
    layout.render_layout(slide, slide_data)
    
    # Export layout configuration
    layout.export_layout('standard_business_layout.json')
    
    print("Modular component system demonstration complete!")