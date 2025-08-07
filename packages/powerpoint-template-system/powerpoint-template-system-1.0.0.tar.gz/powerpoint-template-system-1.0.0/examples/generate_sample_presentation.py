"""
Generate Sample PowerPoint Presentation using the Enhanced Template System

This script demonstrates how to create a professional business presentation
using the new template abstraction system.
"""

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

# Import our template system components
from src.powerpoint_templates.enhanced_business_dsl import BusinessDSLBuilder, BusinessTheme, SlideLayout
from src.powerpoint_templates.business_template_examples import BusinessTemplateExamples


class PowerPointGenerator:
    """Generate PowerPoint presentations using the template system"""
    
    def __init__(self):
        self.presentation = None
        self.slide_number = 0
    
    def create_presentation_from_dsl(self, dsl_data, output_filename="sample_presentation.pptx"):
        """Create PowerPoint presentation from DSL data"""
        
        # Create new presentation
        self.presentation = Presentation()
        self.slide_number = 0
        
        # Set presentation properties
        self.presentation.core_properties.title = dsl_data.title
        self.presentation.core_properties.author = dsl_data.author
        self.presentation.core_properties.subject = dsl_data.subtitle or ""
        
        # Generate slides
        for slide_data in dsl_data.slides:
            self._create_slide(slide_data, dsl_data)
        
        # Save presentation
        self.presentation.save(output_filename)
        return output_filename
    
    def _create_slide(self, slide_data, presentation_data):
        """Create individual slide based on slide data"""
        self.slide_number += 1
        
        # Use blank layout for maximum control
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Route to appropriate slide creation method
        slide_type = slide_data.slide_type
        
        if slide_type == "title":
            self._create_title_slide(slide, slide_data, presentation_data)
        elif slide_type == "agenda":
            self._create_agenda_slide(slide, slide_data, presentation_data)
        elif slide_type == "section":
            self._create_section_slide(slide, slide_data, presentation_data)
        elif slide_type == "content":
            self._create_content_slide(slide, slide_data, presentation_data)
        elif slide_type == "thank_you":
            self._create_thank_you_slide(slide, slide_data, presentation_data)
        else:
            self._create_default_slide(slide, slide_data, presentation_data)
    
    def _create_title_slide(self, slide, slide_data, presentation_data):
        """Create title slide"""
        # Background color
        self._set_slide_background(slide, "#1f4e79")
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = presentation_data.title
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Calibri"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if presentation_data.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.8), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = presentation_data.subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.name = "Calibri"
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Author and company
        author_text = f"{presentation_data.author}"
        if presentation_data.company:
            author_text += f" | {presentation_data.company}"
        if presentation_data.date:
            author_text += f" | {presentation_data.date}"
        
        author_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5), Inches(8), Inches(0.8)
        )
        author_frame = author_box.text_frame
        author_frame.text = author_text
        author_para = author_frame.paragraphs[0]
        author_para.font.name = "Calibri"
        author_para.font.size = Pt(16)
        author_para.font.color.rgb = RGBColor(180, 180, 180)
        author_para.alignment = PP_ALIGN.CENTER
    
    def _create_agenda_slide(self, slide, slide_data, presentation_data):
        """Create agenda slide"""
        # Add header
        self._add_slide_header(slide, slide_data.header.title if slide_data.header else "Agenda", presentation_data)
        
        # Add agenda items
        if slide_data.content and 'items' in slide_data.content.data:
            items = slide_data.content.data['items']
            
            agenda_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2), Inches(7), Inches(4)
            )
            agenda_frame = agenda_box.text_frame
            agenda_frame.clear()
            
            for i, item in enumerate(items):
                if i == 0:
                    para = agenda_frame.paragraphs[0]
                else:
                    para = agenda_frame.add_paragraph()
                
                para.text = f"{i+1}. {item}"
                para.font.name = "Calibri"
                para.font.size = Pt(20)
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.space_after = Pt(12)
        
        # Add footer
        self._add_slide_footer(slide, presentation_data)
    
    def _create_section_slide(self, slide, slide_data, presentation_data):
        """Create section divider slide"""
        # Background color (lighter than title slide)
        self._set_slide_background(slide, "#2c5282")
        
        # Section title
        if slide_data.content and 'title' in slide_data.content.data:
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data.content.data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.name = "Calibri"
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
        
        # Section subtitle
        if slide_data.content and 'subtitle' in slide_data.content.data and slide_data.content.data['subtitle']:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data.content.data['subtitle']
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.name = "Calibri"
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, slide, slide_data, presentation_data):
        """Create content slide"""
        # Add header
        header_title = slide_data.header.title if slide_data.header else "Content"
        self._add_slide_header(slide, header_title, presentation_data)
        
        # Add content based on type
        if slide_data.content:
            content_type = slide_data.content.content_type
            
            if content_type == "bullet_list":
                self._add_bullet_list_content(slide, slide_data.content.data)
            elif content_type == "text":
                self._add_text_content(slide, slide_data.content.data)
            elif content_type == "chart":
                self._add_chart_placeholder(slide, slide_data.content.data)
            else:
                self._add_text_content(slide, {"text": "Content placeholder"})
        
        # Add footer
        self._add_slide_footer(slide, presentation_data)
    
    def _create_thank_you_slide(self, slide, slide_data, presentation_data):
        """Create thank you slide"""
        # Background color
        self._set_slide_background(slide, "#1f4e79")
        
        # Thank you message
        thank_you_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1.5)
        )
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "Thank You"
        thank_you_para = thank_you_frame.paragraphs[0]
        thank_you_para.font.name = "Calibri"
        thank_you_para.font.size = Pt(48)
        thank_you_para.font.bold = True
        thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
        thank_you_para.alignment = PP_ALIGN.CENTER
        
        # Contact information
        if slide_data.content and 'contact_info' in slide_data.content.data:
            contact_info = slide_data.content.data['contact_info']
            contact_text = []
            
            if 'email' in contact_info:
                contact_text.append(f"Email: {contact_info['email']}")
            if 'phone' in contact_info:
                contact_text.append(f"Phone: {contact_info['phone']}")
            if 'website' in contact_info:
                contact_text.append(f"Website: {contact_info['website']}")
            
            if contact_text:
                contact_box = slide.shapes.add_textbox(
                    Inches(1), Inches(4.5), Inches(8), Inches(1.5)
                )
                contact_frame = contact_box.text_frame
                contact_frame.text = "\n".join(contact_text)
                contact_para = contact_frame.paragraphs[0]
                contact_para.font.name = "Calibri"
                contact_para.font.size = Pt(18)
                contact_para.font.color.rgb = RGBColor(200, 200, 200)
                contact_para.alignment = PP_ALIGN.CENTER
    
    def _create_default_slide(self, slide, slide_data, presentation_data):
        """Create default slide layout"""
        self._add_slide_header(slide, "Default Slide", presentation_data)
        self._add_text_content(slide, {"text": "Default slide content"})
        self._add_slide_footer(slide, presentation_data)
    
    def _add_slide_header(self, slide, title, presentation_data):
        """Add header to slide"""
        # Header background
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.2)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = RGBColor(31, 78, 121)  # Corporate blue
        header_bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Calibri"
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.LEFT
        
        # Slide number
        slide_num_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(0.2), Inches(1), Inches(0.8)
        )
        slide_num_frame = slide_num_box.text_frame
        slide_num_frame.text = str(self.slide_number)
        slide_num_para = slide_num_frame.paragraphs[0]
        slide_num_para.font.name = "Calibri"
        slide_num_para.font.size = Pt(16)
        slide_num_para.font.color.rgb = RGBColor(255, 255, 255)
        slide_num_para.alignment = PP_ALIGN.RIGHT
    
    def _add_slide_footer(self, slide, presentation_data):
        """Add footer to slide"""
        # Footer background
        footer_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8), Inches(10), Inches(0.7)
        )
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = RGBColor(240, 240, 240)
        footer_bg.line.color.rgb = RGBColor(200, 200, 200)
        
        # Footer text
        footer_text = []
        if presentation_data.company:
            footer_text.append(presentation_data.company)
        footer_text.append("Confidential")
        footer_text.append(f"Page {self.slide_number}")
        
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.9), Inches(9), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = " | ".join(footer_text)
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.name = "Calibri"
        footer_para.font.size = Pt(10)
        footer_para.font.color.rgb = RGBColor(100, 100, 100)
        footer_para.alignment = PP_ALIGN.CENTER
    
    def _add_bullet_list_content(self, slide, content_data):
        """Add bullet list content to slide"""
        if 'items' not in content_data:
            return
        
        items = content_data['items']
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.clear()
        
        for i, item in enumerate(items):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            if isinstance(item, dict):
                text = item.get('text', '')
                level = item.get('level', 0)
            else:
                text = str(item)
                level = 0
            
            para.text = text
            para.level = min(level, 2)  # Max 3 levels (0, 1, 2)
            para.font.name = "Calibri"
            para.font.size = Pt(max(12, 18 - level * 2))
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.space_after = Pt(6)
    
    def _add_text_content(self, slide, content_data):
        """Add text content to slide"""
        text = content_data.get('text', 'No content provided')
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.text = text
        content_para = content_frame.paragraphs[0]
        content_para.font.name = "Calibri"
        content_para.font.size = Pt(16)
        content_para.font.color.rgb = RGBColor(0, 0, 0)
        content_para.alignment = PP_ALIGN.LEFT
    
    def _add_chart_placeholder(self, slide, content_data):
        """Add chart placeholder to slide"""
        # Chart background
        chart_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(2), Inches(7), Inches(3.5)
        )
        chart_bg.fill.solid()
        chart_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        chart_bg.line.color.rgb = RGBColor(200, 200, 200)
        
        # Chart title
        chart_title = content_data.get('title', 'Chart Placeholder')
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.2), Inches(7), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = chart_title
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Calibri"
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Chart data placeholder
        if 'data' in content_data and 'labels' in content_data:
            data_text = f"Data: {content_data['data']}\nLabels: {content_data['labels']}"
            data_box = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(6), Inches(2)
            )
            data_frame = data_box.text_frame
            data_frame.text = data_text
            data_para = data_frame.paragraphs[0]
            data_para.font.name = "Calibri"
            data_para.font.size = Pt(12)
            data_para.font.color.rgb = RGBColor(100, 100, 100)
            data_para.alignment = PP_ALIGN.CENTER
    
    def _set_slide_background(self, slide, color_hex):
        """Set slide background color"""
        # Convert hex to RGB
        color_hex = color_hex.lstrip('#')
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        
        # Add background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(r, g, b)
        bg_shape.line.fill.background()
        
        # Send to back
        bg_shape.element.getparent().remove(bg_shape.element)
        slide.shapes._spTree.insert(2, bg_shape.element)


def generate_sample_presentations():
    """Generate sample presentations using the template system"""
    
    generator = PowerPointGenerator()
    
    # 1. Generate Quarterly Business Review
    print("Generating Quarterly Business Review presentation...")
    qbr_dsl = BusinessTemplateExamples.create_quarterly_business_review()
    qbr_file = generator.create_presentation_from_dsl(qbr_dsl, "quarterly_business_review.pptx")
    print(f"✅ Created: {qbr_file}")
    
    # 2. Generate Sales Pitch
    print("Generating Sales Pitch presentation...")
    sales_dsl = BusinessTemplateExamples.create_sales_pitch_presentation()
    sales_file = generator.create_presentation_from_dsl(sales_dsl, "sales_pitch_presentation.pptx")
    print(f"✅ Created: {sales_file}")
    
    # 3. Generate Investor Pitch Deck
    print("Generating Investor Pitch Deck...")
    investor_dsl = BusinessTemplateExamples.create_investor_pitch_deck()
    investor_file = generator.create_presentation_from_dsl(investor_dsl, "investor_pitch_deck.pptx")
    print(f"✅ Created: {investor_file}")
    
    # 4. Generate Project Status Report
    print("Generating Project Status Report...")
    project_dsl = BusinessTemplateExamples.create_project_status_report()
    project_file = generator.create_presentation_from_dsl(project_dsl, "project_status_report.pptx")
    print(f"✅ Created: {project_file}")
    
    # 5. Generate Custom Demo Presentation
    print("Generating Custom Demo presentation...")
    demo_dsl = (BusinessDSLBuilder()
                .set_metadata(
                    title="Template System Demo",
                    subtitle="PowerPoint Template Abstraction Showcase",
                    author="AI Assistant",
                    company="@cpro",
                    date=datetime.now().strftime("%B %d, %Y")
                )
                .set_theme(BusinessTheme.MODERN_MINIMAL)
                .add_title_slide()
                .add_agenda_slide("agenda", [
                    "Template System Overview",
                    "Key Features",
                    "Business Templates",
                    "Component Architecture",
                    "Implementation Benefits",
                    "Next Steps"
                ])
                .add_section_divider("overview_section", "Template System Overview", "High-Level Abstraction for Business Presentations")
                .add_content_slide(
                    "key_features",
                    "Key Features of the Template System",
                    "bullet_list",
                    {
                        "items": [
                            {"text": "Business-Focused Templates", "level": 0},
                            {"text": "Executive Summary, Sales Pitch, Investor Deck", "level": 1},
                            {"text": "Professional styling and consistent branding", "level": 1},
                            {"text": "Modular Component System", "level": 0},
                            {"text": "Reusable header, content, and footer components", "level": 1},
                            {"text": "Flexible layouts and positioning", "level": 1},
                            {"text": "Enhanced DSL Structure", "level": 0},
                            {"text": "Business-oriented presentation definition", "level": 1},
                            {"text": "Improved validation and error handling", "level": 1},
                            {"text": "Configuration-Driven Approach", "level": 0},
                            {"text": "External theme and style configuration", "level": 1},
                            {"text": "Easy customization without code changes", "level": 1}
                        ]
                    }
                )
                .add_content_slide(
                    "benefits",
                    "Implementation Benefits",
                    "bullet_list",
                    {
                        "items": [
                            "🚀 60-70% faster presentation creation",
                            "🎨 Consistent professional branding",
                            "🔧 Better code maintainability",
                            "📊 Business-focused templates",
                            "⚙️ Easy customization and theming",
                            "🔄 Seamless integration with existing code",
                            "📚 Comprehensive documentation",
                            "✅ Improved validation and error handling"
                        ]
                    }
                )
                .add_section_divider("next_steps_section", "Next Steps", "Implementation and Migration Strategy")
                .add_content_slide(
                    "implementation_plan",
                    "Recommended Implementation Plan",
                    "bullet_list",
                    {
                        "items": [
                            {"text": "Phase 1: Foundation (Weeks 1-2)", "level": 0},
                            {"text": "Implement core template system", "level": 1},
                            {"text": "Set up component framework", "level": 1},
                            {"text": "Phase 2: DSL Enhancement (Weeks 3-4)", "level": 0},
                            {"text": "Deploy enhanced DSL structure", "level": 1},
                            {"text": "Create integration layer", "level": 1},
                            {"text": "Phase 3: Business Templates (Weeks 5-6)", "level": 0},
                            {"text": "Implement business presentation templates", "level": 1},
                            {"text": "Comprehensive testing and optimization", "level": 1}
                        ]
                    }
                )
                .add_thank_you_slide(
                    contact_info={
                        "email": "research@cpro.com",
                        "demo": "Template system demonstration complete"
                    }
                )
                .build())
    
    demo_file = generator.create_presentation_from_dsl(demo_dsl, "template_system_demo.pptx")
    print(f"✅ Created: {demo_file}")
    
    return [qbr_file, sales_file, investor_file, project_file, demo_file]


if __name__ == "__main__":
    print("=== PowerPoint Template System - Presentation Generation ===\n")
    
    try:
        # Generate all sample presentations
        generated_files = generate_sample_presentations()
        
        print(f"\n🎉 Successfully generated {len(generated_files)} PowerPoint presentations!")
        print("\nGenerated files:")
        for file in generated_files:
            file_size = os.path.getsize(file) / 1024  # Size in KB
            print(f"  📄 {file} ({file_size:.1f} KB)")
        
        print("\n✨ All presentations demonstrate the enhanced template system features:")
        print("  • Professional business styling")
        print("  • Consistent branding and layouts")
        print("  • Modular component architecture")
        print("  • Business-focused content structure")
        print("  • Enhanced DSL capabilities")
        
    except Exception as e:
        print(f"❌ Error generating presentations: {e}")
        import traceback
        traceback.print_exc()