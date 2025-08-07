"""
Enhanced Visual PowerPoint Generator with Charts, Hero Art, and 16:9 Aspect Ratio

This module creates visually appealing PowerPoint presentations with:
- Proper 16:9 widescreen aspect ratio
- Hero artwork and visual elements
- Actual charts and data visualizations
- Professional graphics and layouts
"""

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.dml import MSO_THEME_COLOR
import os
from datetime import datetime
import math

# Import our template system components
from .enhanced_business_dsl import BusinessDSLBuilder, BusinessTheme, SlideLayout
from .business_template_examples import BusinessTemplateExamples


class EnhancedVisualGenerator:
    """Enhanced PowerPoint generator with visual improvements"""
    
    def __init__(self):
        self.presentation = None
        self.slide_number = 0
        # 16:9 aspect ratio dimensions (in inches)
        self.slide_width = 13.33
        self.slide_height = 7.5
        
        # Color palettes for different themes
        self.color_palettes = {
            'corporate_blue': {
                'primary': RGBColor(31, 78, 121),      # #1f4e79
                'secondary': RGBColor(112, 173, 71),    # #70ad47
                'accent': RGBColor(197, 90, 17),        # #c55a11
                'light': RGBColor(240, 248, 255),       # #f0f8ff
                'dark': RGBColor(25, 25, 25),           # #191919
                'gradient_start': RGBColor(31, 78, 121),
                'gradient_end': RGBColor(44, 82, 130)
            },
            'modern_minimal': {
                'primary': RGBColor(44, 62, 80),        # #2c3e50
                'secondary': RGBColor(149, 165, 166),   # #95a5a6
                'accent': RGBColor(52, 152, 219),       # #3498db
                'light': RGBColor(250, 250, 250),       # #fafafa
                'dark': RGBColor(44, 62, 80),           # #2c3e50
                'gradient_start': RGBColor(44, 62, 80),
                'gradient_end': RGBColor(52, 73, 94)
            },
            'startup_vibrant': {
                'primary': RGBColor(231, 76, 60),       # #e74c3c
                'secondary': RGBColor(241, 196, 15),    # #f1c40f
                'accent': RGBColor(155, 89, 182),       # #9b59b6
                'light': RGBColor(255, 250, 240),       # #fffaf0
                'dark': RGBColor(44, 62, 80),           # #2c3e50
                'gradient_start': RGBColor(231, 76, 60),
                'gradient_end': RGBColor(192, 57, 43)
            }
        }
    
    def create_presentation_from_dsl(self, dsl_data, output_filename="enhanced_presentation.pptx"):
        """Create enhanced PowerPoint presentation from DSL data"""
        
        # Create new presentation with 16:9 aspect ratio
        self.presentation = Presentation()
        
        # Set slide size to 16:9 widescreen
        self.presentation.slide_width = Inches(self.slide_width)
        self.presentation.slide_height = Inches(self.slide_height)
        
        # Set presentation properties
        self.presentation.core_properties.title = dsl_data.title
        self.presentation.core_properties.author = dsl_data.author
        self.presentation.core_properties.subject = dsl_data.subtitle or ""
        
        # Get color palette for theme
        theme_name = dsl_data.theme.value if hasattr(dsl_data.theme, 'value') else 'corporate_blue'
        self.current_palette = self.color_palettes.get(theme_name, self.color_palettes['corporate_blue'])
        
        self.slide_number = 0
        
        # Generate slides with enhanced visuals
        for slide_data in dsl_data.slides:
            self._create_enhanced_slide(slide_data, dsl_data)
        
        # Save presentation
        self.presentation.save(output_filename)
        return output_filename
    
    def _create_enhanced_slide(self, slide_data, presentation_data):
        """Create enhanced slide with improved visuals"""
        self.slide_number += 1
        
        # Use blank layout for maximum control
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Route to appropriate enhanced slide creation method
        slide_type = slide_data.slide_type
        
        if slide_type == "title":
            self._create_enhanced_title_slide(slide, slide_data, presentation_data)
        elif slide_type == "agenda":
            self._create_enhanced_agenda_slide(slide, slide_data, presentation_data)
        elif slide_type == "section":
            self._create_enhanced_section_slide(slide, slide_data, presentation_data)
        elif slide_type == "content":
            self._create_enhanced_content_slide(slide, slide_data, presentation_data)
        elif slide_type == "thank_you":
            self._create_enhanced_thank_you_slide(slide, slide_data, presentation_data)
        else:
            self._create_enhanced_default_slide(slide, slide_data, presentation_data)
    
    def _create_enhanced_title_slide(self, slide, slide_data, presentation_data):
        """Create enhanced title slide with hero artwork"""
        # Create gradient background
        self._add_gradient_background(slide, self.current_palette['gradient_start'], self.current_palette['gradient_end'])
        
        # Add hero geometric shapes
        self._add_hero_shapes(slide)
        
        # Main title with enhanced styling
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), Inches(10.33), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = presentation_data.title
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Segoe UI"
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add text shadow effect
        title_para.font.shadow = True
        
        # Subtitle with modern styling
        if presentation_data.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.2), Inches(10.33), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = presentation_data.subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.name = "Segoe UI Light"
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(220, 220, 220)
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Author information with modern card design
        self._add_author_card(slide, presentation_data)
    
    def _create_enhanced_agenda_slide(self, slide, slide_data, presentation_data):
        """Create enhanced agenda slide with visual elements"""
        # Add subtle background pattern
        self._add_subtle_background(slide)
        
        # Enhanced header
        self._add_enhanced_header(slide, slide_data.header.title if slide_data.header else "Agenda", presentation_data)
        
        # Add agenda items with enhanced styling
        if slide_data.content and 'items' in slide_data.content.data:
            items = slide_data.content.data['items']
            self._add_enhanced_agenda_items(slide, items)
        
        # Enhanced footer
        self._add_enhanced_footer(slide, presentation_data)
    
    def _create_enhanced_section_slide(self, slide, slide_data, presentation_data):
        """Create enhanced section divider with dramatic visuals"""
        # Create diagonal gradient background
        self._add_diagonal_gradient_background(slide)
        
        # Add decorative elements
        self._add_section_decorations(slide)
        
        # Section title with enhanced typography
        if slide_data.content and 'title' in slide_data.content.data:
            title_box = slide.shapes.add_textbox(
                Inches(2), Inches(2.5), Inches(9.33), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data.content.data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.name = "Segoe UI"
            title_para.font.size = Pt(42)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
        
        # Section subtitle
        if slide_data.content and 'subtitle' in slide_data.content.data and slide_data.content.data['subtitle']:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.5), Inches(9.33), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data.content.data['subtitle']
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.name = "Segoe UI Light"
            subtitle_para.font.size = Pt(22)
            subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def _create_enhanced_content_slide(self, slide, slide_data, presentation_data):
        """Create enhanced content slide with charts and visuals"""
        # Add subtle background
        self._add_subtle_background(slide)
        
        # Enhanced header
        header_title = slide_data.header.title if slide_data.header else "Content"
        self._add_enhanced_header(slide, header_title, presentation_data)
        
        # Add content based on type with enhancements
        if slide_data.content:
            content_type = slide_data.content.content_type
            
            if content_type == "bullet_list":
                self._add_enhanced_bullet_list(slide, slide_data.content.data)
            elif content_type == "text":
                self._add_enhanced_text_content(slide, slide_data.content.data)
            elif content_type == "chart":
                self._add_actual_chart(slide, slide_data.content.data)
            elif content_type == "image_with_text":
                self._add_image_with_text(slide, slide_data.content.data)
            elif content_type == "image_focus":
                self._add_image_focus(slide, slide_data.content.data)
            elif content_type == "image_grid":
                self._add_image_grid(slide, slide_data.content.data)
            elif content_type == "icon_grid":
                self._add_icon_grid(slide, slide_data.content.data)
            elif content_type == "mixed_content":
                self._add_mixed_content(slide, slide_data.content.data)
            elif content_type == "card_grid":
                self._add_card_grid(slide, slide_data.content.data)
            else:
                self._add_enhanced_text_content(slide, {"text": "Enhanced content placeholder"})
        
        # Enhanced footer
        self._add_enhanced_footer(slide, presentation_data)
    
    def _create_enhanced_thank_you_slide(self, slide, slide_data, presentation_data):
        """Create enhanced thank you slide"""
        # Create radial gradient background
        self._add_radial_gradient_background(slide)
        
        # Add celebratory elements
        self._add_thank_you_decorations(slide)
        
        # Thank you message with enhanced styling
        thank_you_box = slide.shapes.add_textbox(
            Inches(2), Inches(2), Inches(9.33), Inches(1.5)
        )
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "Thank You"
        thank_you_para = thank_you_frame.paragraphs[0]
        thank_you_para.font.name = "Segoe UI"
        thank_you_para.font.size = Pt(54)
        thank_you_para.font.bold = True
        thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
        thank_you_para.alignment = PP_ALIGN.CENTER
        
        # Contact information card
        if slide_data.content and 'contact_info' in slide_data.content.data:
            self._add_contact_card(slide, slide_data.content.data['contact_info'])
    
    def _add_gradient_background(self, slide, start_color, end_color):
        """Add gradient background to slide"""
        # Create background rectangle
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.slide_width), Inches(self.slide_height)
        )
        
        # Set gradient fill
        fill = bg_shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = start_color
        fill.gradient_stops[1].color.rgb = end_color
        
        # Remove border
        bg_shape.line.fill.background()
        
        # Send to back
        self._send_to_back(slide, bg_shape)
    
    def _add_diagonal_gradient_background(self, slide):
        """Add diagonal gradient background"""
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.slide_width), Inches(self.slide_height)
        )
        
        fill = bg_shape.fill
        fill.gradient()
        fill.gradient_angle = 45  # Diagonal gradient
        fill.gradient_stops[0].color.rgb = self.current_palette['primary']
        fill.gradient_stops[1].color.rgb = self.current_palette['secondary']
        
        bg_shape.line.fill.background()
        self._send_to_back(slide, bg_shape)
    
    def _add_radial_gradient_background(self, slide):
        """Add radial gradient background"""
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.slide_width), Inches(self.slide_height)
        )
        
        fill = bg_shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.current_palette['primary']
        fill.gradient_stops[1].color.rgb = self.current_palette['dark']
        
        bg_shape.line.fill.background()
        self._send_to_back(slide, bg_shape)
    
    def _add_subtle_background(self, slide):
        """Add subtle background pattern"""
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.slide_width), Inches(self.slide_height)
        )
        
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        bg_shape.line.fill.background()
        self._send_to_back(slide, bg_shape)
        
        # Add subtle pattern elements
        for i in range(5):
            for j in range(3):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(2 + i * 2), Inches(1 + j * 2), Inches(0.05), Inches(0.05)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = RGBColor(230, 230, 230)
                dot.line.fill.background()
    
    def _add_hero_shapes(self, slide):
        """Add hero geometric shapes to title slide"""
        # Large circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(10), Inches(0.5), Inches(4), Inches(4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.current_palette['accent']
        circle.fill.transparency = 0.3
        circle.line.fill.background()
        
        # Triangle
        triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(0.5), Inches(4), Inches(2), Inches(2)
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = self.current_palette['secondary']
        triangle.fill.transparency = 0.4
        triangle.line.fill.background()
        
        # Rectangle
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(11.5), Inches(5), Inches(1.5), Inches(1.5)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = self.current_palette['light']
        rect.fill.transparency = 0.2
        rect.line.fill.background()
    
    def _add_section_decorations(self, slide):
        """Add decorative elements to section slides"""
        # Add flowing lines
        for i in range(3):
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + i * 0.3), Inches(6), Inches(11.33 - i * 0.6), Inches(0.1)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            line_shape.fill.transparency = 0.7 - i * 0.2
            line_shape.line.fill.background()
    
    def _add_thank_you_decorations(self, slide):
        """Add celebratory elements to thank you slide"""
        # Add sparkle elements
        sparkle_positions = [
            (2, 1.5), (11, 2), (1.5, 5.5), (10.5, 6), (6.5, 1), (3.5, 6.5)
        ]
        
        for x, y in sparkle_positions:
            sparkle = slide.shapes.add_shape(
                MSO_SHAPE.STAR_4_POINT,
                Inches(x), Inches(y), Inches(0.3), Inches(0.3)
            )
            sparkle.fill.solid()
            sparkle.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold
            sparkle.line.fill.background()
    
    def _add_enhanced_header(self, slide, title, presentation_data):
        """Add enhanced header with modern styling"""
        # Header background with gradient
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.slide_width), Inches(1.2)
        )
        
        fill = header_bg.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.current_palette['primary']
        fill.gradient_stops[1].color.rgb = self.current_palette['gradient_end']
        header_bg.line.fill.background()
        
        # Title with enhanced typography
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.2), Inches(10), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Segoe UI"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.LEFT
        
        # Slide number with modern styling
        slide_num_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(11.8), Inches(0.3), Inches(0.6), Inches(0.6)
        )
        slide_num_bg.fill.solid()
        slide_num_bg.fill.fore_color.rgb = self.current_palette['accent']
        slide_num_bg.line.fill.background()
        
        slide_num_box = slide.shapes.add_textbox(
            Inches(11.8), Inches(0.3), Inches(0.6), Inches(0.6)
        )
        slide_num_frame = slide_num_box.text_frame
        slide_num_frame.text = str(self.slide_number)
        slide_num_para = slide_num_frame.paragraphs[0]
        slide_num_para.font.name = "Segoe UI"
        slide_num_para.font.size = Pt(16)
        slide_num_para.font.bold = True
        slide_num_para.font.color.rgb = RGBColor(255, 255, 255)
        slide_num_para.alignment = PP_ALIGN.CENTER
        slide_num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _add_enhanced_footer(self, slide, presentation_data):
        """Add enhanced footer with modern styling"""
        # Footer background
        footer_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8), Inches(self.slide_width), Inches(0.7)
        )
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        footer_bg.line.color.rgb = RGBColor(220, 220, 220)
        
        # Footer text with enhanced styling
        footer_text = []
        if presentation_data.company:
            footer_text.append(presentation_data.company)
        footer_text.append("Confidential")
        footer_text.append(f"Page {self.slide_number}")
        
        footer_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.9), Inches(11.73), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = " | ".join(footer_text)
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.name = "Segoe UI"
        footer_para.font.size = Pt(11)
        footer_para.font.color.rgb = RGBColor(100, 100, 100)
        footer_para.alignment = PP_ALIGN.CENTER
    
    def _add_enhanced_bullet_list(self, slide, content_data):
        """Add enhanced bullet list with modern styling"""
        if 'items' not in content_data:
            return
        
        items = content_data['items']
        
        # Content background card
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.5), Inches(11.33), Inches(4.8)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_bg.line.color.rgb = RGBColor(230, 230, 230)
        content_bg.shadow.inherit = False
        
        # Content text
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.2)
        )
        content_frame = content_box.text_frame
        content_frame.clear()
        
        for i, item in enumerate(items):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            if isinstance(item, dict):
                text = item.get('text', '')
                level = item.get('level', 0)
            else:
                text = str(item)
                level = 0
            
            para.text = text
            para.level = min(level, 2)
            para.font.name = "Segoe UI"
            para.font.size = Pt(max(14, 20 - level * 2))
            para.font.color.rgb = self.current_palette['dark']
            para.space_after = Pt(8)
            
            # Add colored bullet for level 0
            if level == 0:
                para.font.bold = True
    
    def _add_enhanced_text_content(self, slide, content_data):
        """Add enhanced text content with modern styling"""
        text = content_data.get('text', 'No content provided')
        
        # Content background card
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.5), Inches(11.33), Inches(4.8)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_bg.line.color.rgb = RGBColor(230, 230, 230)
        
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.2)
        )
        content_frame = content_box.text_frame
        content_frame.text = text
        content_para = content_frame.paragraphs[0]
        content_para.font.name = "Segoe UI"
        content_para.font.size = Pt(18)
        content_para.font.color.rgb = self.current_palette['dark']
        content_para.alignment = PP_ALIGN.LEFT
        content_para.line_spacing = 1.2
    
    def _add_actual_chart(self, slide, content_data):
        """Add actual chart instead of placeholder"""
        chart_type = content_data.get('chart_type', 'line')
        chart_data = content_data.get('data', [100, 120, 135, 150])
        chart_labels = content_data.get('labels', ['Q1', 'Q2', 'Q3', 'Q4'])
        chart_title = content_data.get('title', 'Chart')
        
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_labels
        chart_data_obj.add_series('Series 1', chart_data)
        
        # Determine chart type
        if chart_type == 'line':
            chart_type_enum = XL_CHART_TYPE.LINE
        elif chart_type == 'bar':
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type == 'pie':
            chart_type_enum = XL_CHART_TYPE.PIE
        else:
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        
        # Add chart to slide
        chart_frame = slide.shapes.add_chart(
            chart_type_enum, 
            Inches(2), Inches(2), Inches(9.33), Inches(4),
            chart_data_obj
        )
        
        chart = chart_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        
        # Style the chart
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.current_palette['primary']
        
        # Style chart series
        if hasattr(chart.series[0], 'format'):
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = self.current_palette['primary']
    
    def _add_actual_chart_with_position(self, slide, content_data, left, top, width, height):
        """Add actual chart with custom positioning"""
        chart_type = content_data.get('chart_type', 'bar')
        data = content_data.get('data', [])
        labels = content_data.get('labels', [])
        chart_title = content_data.get('title', 'Chart')
        
        # Create chart with custom position
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series('Series 1', data)
        
        # Add chart to slide with custom positioning
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED if chart_type == 'bar' else XL_CHART_TYPE.LINE,
            Inches(left), Inches(top), Inches(width), Inches(height), chart_data
        ).chart
        
        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        
        # Style the chart
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.current_palette['primary']
        
        # Style chart series
        if hasattr(chart.series[0], 'format'):
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = self.current_palette['primary']
    
    def _add_enhanced_agenda_items(self, slide, items):
        """Add enhanced agenda items with visual elements"""
        # Background card for agenda
        agenda_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2), Inches(2), Inches(9.33), Inches(4)
        )
        agenda_bg.fill.solid()
        agenda_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        agenda_bg.line.color.rgb = RGBColor(230, 230, 230)
        
        # Add agenda items with numbers and styling
        for i, item in enumerate(items):
            y_pos = 2.3 + i * 0.6
            
            # Number circle
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(2.5), Inches(y_pos), Inches(0.4), Inches(0.4)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = self.current_palette['primary']
            num_circle.line.fill.background()
            
            # Number text
            num_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(y_pos), Inches(0.4), Inches(0.4)
            )
            num_frame = num_box.text_frame
            num_frame.text = str(i + 1)
            num_para = num_frame.paragraphs[0]
            num_para.font.name = "Segoe UI"
            num_para.font.size = Pt(14)
            num_para.font.bold = True
            num_para.font.color.rgb = RGBColor(255, 255, 255)
            num_para.alignment = PP_ALIGN.CENTER
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Item text
            item_box = slide.shapes.add_textbox(
                Inches(3.2), Inches(y_pos), Inches(7.5), Inches(0.4)
            )
            item_frame = item_box.text_frame
            item_frame.text = item
            item_para = item_frame.paragraphs[0]
            item_para.font.name = "Segoe UI"
            item_para.font.size = Pt(18)
            item_para.font.color.rgb = self.current_palette['dark']
            item_para.alignment = PP_ALIGN.LEFT
            item_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _add_author_card(self, slide, presentation_data):
        """Add modern author information card"""
        # Author card background
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4), Inches(5.5), Inches(5.33), Inches(1.2)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card_bg.fill.transparency = 0.1
        card_bg.line.color.rgb = RGBColor(255, 255, 255)
        card_bg.line.transparency = 0.3
        
        # Author text
        author_text = f"{presentation_data.author}"
        if presentation_data.company:
            author_text += f" | {presentation_data.company}"
        if presentation_data.date:
            author_text += f"\n{presentation_data.date}"
        
        author_box = slide.shapes.add_textbox(
            Inches(4.3), Inches(5.7), Inches(4.73), Inches(0.8)
        )
        author_frame = author_box.text_frame
        author_frame.text = author_text
        author_para = author_frame.paragraphs[0]
        author_para.font.name = "Segoe UI"
        author_para.font.size = Pt(16)
        author_para.font.color.rgb = RGBColor(240, 240, 240)
        author_para.alignment = PP_ALIGN.CENTER
    
    def _add_contact_card(self, slide, contact_info):
        """Add modern contact information card"""
        # Contact card background
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.5), Inches(4.5), Inches(6.33), Inches(1.8)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card_bg.fill.transparency = 0.1
        card_bg.line.color.rgb = RGBColor(255, 255, 255)
        card_bg.line.transparency = 0.3
        
        # Contact information
        contact_text = []
        if 'email' in contact_info:
            contact_text.append(f"📧 {contact_info['email']}")
        if 'phone' in contact_info:
            contact_text.append(f"📞 {contact_info['phone']}")
        if 'website' in contact_info:
            contact_text.append(f"🌐 {contact_info['website']}")
        
        if contact_text:
            contact_box = slide.shapes.add_textbox(
                Inches(3.8), Inches(4.8), Inches(5.73), Inches(1.2)
            )
            contact_frame = contact_box.text_frame
            contact_frame.text = "\n".join(contact_text)
            contact_para = contact_frame.paragraphs[0]
            contact_para.font.name = "Segoe UI"
            contact_para.font.size = Pt(16)
            contact_para.font.color.rgb = RGBColor(220, 220, 220)
            contact_para.alignment = PP_ALIGN.CENTER
            contact_para.line_spacing = 1.3
    
    def _add_image_with_text(self, slide, content_data):
        """Add image with text content side by side"""
        image_path = content_data.get('image_path', '')
        image_position = content_data.get('image_position', 'left')
        text_content = content_data.get('text_content', [])
        image_width = content_data.get('image_width', 3.0)
        image_height = content_data.get('image_height', 2.5)
        
        if image_path and os.path.exists(image_path):
            # Add image
            if image_position == 'left':
                img_left = 1.0
                text_left = 1.0 + image_width + 0.5
            else:  # right
                img_left = self.slide_width - image_width - 1.0
                text_left = 1.0
            
            # Insert image
            slide.shapes.add_picture(image_path, Inches(img_left), Inches(2.5), 
                                   Inches(image_width), Inches(image_height))
            
            # Add text content
            text_width = self.slide_width - image_width - 2.5
            text_box = slide.shapes.add_textbox(
                Inches(text_left), Inches(2.5), Inches(text_width), Inches(4)
            )
            text_frame = text_box.text_frame
            text_frame.clear()
            
            for item in text_content:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.name = "Segoe UI"
                p.font.size = Pt(16)
                p.font.color.rgb = self.current_palette['dark']
                p.space_after = Pt(12)
        else:
            # Fallback to text only
            self._add_enhanced_text_content(slide, {"text": "Image not found"})
    
    def _add_image_focus(self, slide, content_data):
        """Add image as the main focus with optional caption"""
        image_path = content_data.get('image_path', '')
        image_width = content_data.get('image_width', 6.0)
        image_height = content_data.get('image_height', 4.0)
        caption = content_data.get('caption', '')
        
        if image_path and os.path.exists(image_path):
            # Center the image
            img_left = (self.slide_width - image_width) / 2
            img_top = (self.slide_height - image_height) / 2
            
            # Insert image
            slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), 
                                   Inches(image_width), Inches(image_height))
            
            # Add caption if provided
            if caption:
                caption_box = slide.shapes.add_textbox(
                    Inches(img_left), Inches(img_top + image_height + 0.2), 
                    Inches(image_width), Inches(0.5)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = caption
                caption_para = caption_frame.paragraphs[0]
                caption_para.font.name = "Segoe UI"
                caption_para.font.size = Pt(14)
                caption_para.font.color.rgb = self.current_palette['dark']
                caption_para.alignment = PP_ALIGN.CENTER
        else:
            # Fallback to text
            self._add_enhanced_text_content(slide, {"text": "Image not found"})
    
    def _add_image_grid(self, slide, content_data):
        """Add multiple images in a grid layout"""
        images = content_data.get('images', [])
        layout = content_data.get('layout', '2x2')
        
        if not images:
            self._add_enhanced_text_content(slide, {"text": "No images provided"})
            return
        
        # Calculate grid dimensions
        if layout == '2x2':
            cols, rows = 2, 2
        elif layout == '3x2':
            cols, rows = 3, 2
        else:
            cols, rows = 2, 2
        
        img_width = (self.slide_width - 2.0) / cols
        img_height = 2.0
        
        for i, img_data in enumerate(images[:cols*rows]):
            if isinstance(img_data, dict):
                img_path = img_data.get('path', '')
                caption = img_data.get('caption', '')
            else:
                img_path = img_data
                caption = ''
            
            if img_path and os.path.exists(img_path):
                # Calculate position
                col = i % cols
                row = i // cols
                
                img_left = 1.0 + col * (img_width + 0.2)
                img_top = 1.8 + row * (img_height + 0.3)  # Moved up from 2.5 and reduced spacing from 0.5 to 0.3
                
                # Insert image
                slide.shapes.add_picture(img_path, Inches(img_left), Inches(img_top), 
                                       Inches(img_width), Inches(img_height))
                
                # Add caption
                if caption:
                    caption_box = slide.shapes.add_textbox(
                        Inches(img_left), Inches(img_top + img_height + 0.05), 
                        Inches(img_width), Inches(0.3)
                    )
                    caption_frame = caption_box.text_frame
                    caption_frame.text = caption
                    caption_para = caption_frame.paragraphs[0]
                    caption_para.font.name = "Segoe UI"
                    caption_para.font.size = Pt(12)
                    caption_para.font.color.rgb = self.current_palette['dark']
                    caption_para.alignment = PP_ALIGN.CENTER
    
    def _add_icon_grid(self, slide, content_data):
        """Add icons in a grid layout with titles and descriptions"""
        icons = content_data.get('icons', [])
        layout = content_data.get('layout', '2x2')
        
        if not icons:
            self._add_enhanced_text_content(slide, {"text": "No icons provided"})
            return
        
        # Calculate grid dimensions
        if layout == '2x2':
            cols, rows = 2, 2
        elif layout == '3x2':
            cols, rows = 3, 2
        else:
            cols, rows = 2, 2
        
        icon_size = 1.5
        spacing = 0.5
        
        for i, icon_data in enumerate(icons[:cols*rows]):
            img_path = icon_data.get('path', '')
            title = icon_data.get('title', '')
            description = icon_data.get('description', '')
            
            if img_path and os.path.exists(img_path):
                # Calculate position
                col = i % cols
                row = i // cols
                
                icon_left = 1.0 + col * (icon_size + spacing)
                icon_top = 1.8 + row * (icon_size + 1.2)  # Moved entire grid up from 2.5 to 1.8
                
                # Insert icon
                slide.shapes.add_picture(img_path, Inches(icon_left), Inches(icon_top), 
                                       Inches(icon_size), Inches(icon_size))
                
                # Add title
                if title:
                    title_box = slide.shapes.add_textbox(
                        Inches(icon_left), Inches(icon_top + icon_size + 0.1), 
                        Inches(icon_size), Inches(0.3)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = title
                    title_para = title_frame.paragraphs[0]
                    title_para.font.name = "Segoe UI"
                    title_para.font.size = Pt(14)
                    title_para.font.bold = True
                    title_para.font.color.rgb = self.current_palette['primary']
                    title_para.alignment = PP_ALIGN.CENTER
                
                # Add description
                if description:
                    desc_box = slide.shapes.add_textbox(
                        Inches(icon_left), Inches(icon_top + icon_size + 0.4), 
                        Inches(icon_size), Inches(0.4)
                    )
                    desc_frame = desc_box.text_frame
                    desc_frame.text = description
                    desc_para = desc_frame.paragraphs[0]
                    desc_para.font.name = "Segoe UI"
                    desc_para.font.size = Pt(10)
                    desc_para.font.color.rgb = self.current_palette['dark']
                    desc_para.alignment = PP_ALIGN.CENTER
    
    def _add_mixed_content(self, slide, content_data):
        """Add mixed content (chart + image)"""
        chart_data = content_data.get('chart_data', {})
        image_path = content_data.get('image_path', '')
        image_position = content_data.get('image_position', 'bottom')
        image_width = content_data.get('image_width', 4.0)
        image_height = content_data.get('image_height', 2.0)
        
        # Add chart if provided (smaller size to leave room for image)
        if chart_data:
            # Create a smaller chart area
            chart_left = 1.0
            chart_top = 2.5
            chart_width = self.slide_width - 2.0
            chart_height = 2.2 if image_position == 'bottom' else 4.0  # Reduced height to make room for image
            
            # Add chart with custom positioning
            self._add_actual_chart_with_position(slide, chart_data, chart_left, chart_top, chart_width, chart_height)
        
        # Add image if provided
        if image_path and os.path.exists(image_path):
            if image_position == 'bottom':
                img_left = (self.slide_width - image_width) / 2
                img_top = 4.8  # Moved up from 5.5 to better integrate with main chart
            else:  # right
                img_left = self.slide_width - image_width - 1.0
                img_top = 2.5
            
            slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), 
                                   Inches(image_width), Inches(image_height))
    
    def _add_card_grid(self, slide, content_data):
        """Add article-style cards in a grid layout"""
        cards = content_data.get('cards', [])
        layout = content_data.get('layout', 'horizontal')
        background_color = content_data.get('background_color', '#1f4e79')
        card_spacing = content_data.get('card_spacing', 0.3)
        
        if not cards:
            self._add_enhanced_text_content(slide, {"text": "No cards provided"})
            return
        
        # Set slide background
        self._set_slide_background(slide, background_color)
        
        # Calculate card dimensions and positioning
        num_cards = len(cards)
        if layout == 'horizontal':
            card_width = (self.slide_width - 2.0 - (num_cards - 1) * card_spacing) / num_cards
            card_height = 4.0
            start_x = 1.0
            start_y = 2.0
        else:  # vertical
            card_width = (self.slide_width - 2.0) / 2
            card_height = 3.0
            start_x = 1.0
            start_y = 2.0
        
        # Add slide title
        if content_data.get('title'):
            title_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(0.5), Inches(self.slide_width - 2.0), Inches(1.0)
            )
            title_frame = title_box.text_frame
            title_frame.text = content_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.name = "Segoe UI"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
        
        # Create cards
        for i, card_data in enumerate(cards):
            if layout == 'horizontal':
                card_x = start_x + i * (card_width + card_spacing)
                card_y = start_y
            else:
                row = i // 2
                col = i % 2
                card_x = start_x + col * (card_width + card_spacing)
                card_y = start_y + row * (card_height + 0.5)
            
            self._create_article_card(slide, card_data, card_x, card_y, card_width, card_height)
    
    def _create_article_card(self, slide, card_data, x, y, width, height):
        """Create a single article card with image, category, title, and description"""
        # Get card styling options
        rounded_corners = card_data.get('rounded_corners', True)
        card_color = card_data.get('card_color', '#ffffff')
        gradient = card_data.get('gradient', None)
        shadow = card_data.get('shadow', True)
        border_style = card_data.get('border_style', 'solid')  # solid, dashed, dotted
        border_width = card_data.get('border_width', 1)
        border_color = card_data.get('border_color', '#e0e0e0')
        
        # Convert hex colors to RGB
        card_color = card_color.lstrip('#')
        card_r = int(card_color[0:2], 16)
        card_g = int(card_color[2:4], 16)
        card_b = int(card_color[4:6], 16)
        
        border_color = border_color.lstrip('#')
        border_r = int(border_color[0:2], 16)
        border_g = int(border_color[2:4], 16)
        border_b = int(border_color[4:6], 16)
        
        # Card background with optional rounded corners
        if rounded_corners:
            card_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            # Ensure all corners are rounded by setting the adjustment value
            card_bg.adjustments[0] = 0.1  # This makes the corners more rounded
        else:
            card_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
        
        # Apply gradient or solid fill
        if gradient:
            self._apply_gradient_fill(card_bg, gradient)
        else:
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(card_r, card_g, card_b)
        
        # Apply border styling
        card_bg.line.color.rgb = RGBColor(border_r, border_g, border_b)
        card_bg.line.width = Pt(border_width)
        
        # Apply border style (simplified for compatibility)
        # Note: python-pptx has limited line style support
        card_bg.line.dash_style = MSO_LINE.SOLID
        
        # Add shadow effect
        if shadow:
            self._add_card_shadow(slide, x, y, width, height, rounded_corners)
        
        # Calculate image dimensions
        image_height = height * 0.4
        image_width = width
        image_x = x
        image_y = y
        
        # Add image if provided
        image_path = card_data.get('image_path', '')
        image_rounded = card_data.get('image_rounded', None)  # New option for image rounding
        
        if image_path and os.path.exists(image_path):
            try:
                # Determine if image should be rounded (use card setting if not specified)
                should_round_image = image_rounded if image_rounded is not None else rounded_corners
                
                if should_round_image:
                    # Create a rounded rectangle placeholder for rounded images
                    placeholder = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(image_x), Inches(image_y), Inches(image_width), Inches(image_height)
                    )
                    placeholder.adjustments[0] = 0.1  # Match the card's corner rounding
                    
                    # Set placeholder color based on image name
                    img_name = os.path.basename(image_path).lower()
                    if 'product' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(100, 150, 100)  # Green
                    elif 'process' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(150, 100, 50)  # Brown
                    elif 'team' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(100, 100, 150)  # Purple
                    elif 'idea' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(200, 150, 100)  # Light brown
                    elif 'target' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(150, 100, 150)  # Purple
                    elif 'results' in img_name or 'summary' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(80, 80, 80)  # Dark gray
                    elif 'chart' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(70, 130, 180)  # Blue
                    elif 'growth' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(100, 150, 100)  # Green
                    elif 'company' in img_name or 'logo' in img_name:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(120, 120, 120)  # Gray
                    else:
                        placeholder.fill.solid()
                        placeholder.fill.fore_color.rgb = RGBColor(120, 120, 120)  # Gray
                    
                    # Add text overlay
                    text_shape = slide.shapes.add_textbox(
                        Inches(image_x), Inches(image_y), Inches(image_width), Inches(image_height)
                    )
                    text_frame = text_shape.text_frame
                    text_frame.text = os.path.basename(image_path).replace('.png', '').replace('_', ' ').title()
                    text_para = text_frame.paragraphs[0]
                    text_para.font.name = "Segoe UI"
                    text_para.font.size = Pt(12)
                    text_para.font.bold = True
                    text_para.font.color.rgb = RGBColor(255, 255, 255)
                    text_para.alignment = PP_ALIGN.CENTER
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                else:
                    # For square corners, use the actual image
                    slide.shapes.add_picture(image_path, Inches(image_x), Inches(image_y), 
                                           Inches(image_width), Inches(image_height))
                    
            except:
                # Fallback if image fails to load - create a colored placeholder
                should_round_image = image_rounded if image_rounded is not None else rounded_corners
                
                if should_round_image:
                    placeholder = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(image_x), Inches(image_y), Inches(image_width), Inches(image_height)
                    )
                    placeholder.adjustments[0] = 0.1
                else:
                    placeholder = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(image_x), Inches(image_y), Inches(image_width), Inches(image_height)
                    )
                
                # Set placeholder color based on image name
                img_name = os.path.basename(image_path).lower()
                if 'product' in img_name:
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(100, 150, 100)  # Green
                elif 'process' in img_name:
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(150, 100, 50)  # Brown
                elif 'team' in img_name:
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(100, 100, 150)  # Purple
                else:
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(120, 120, 120)  # Gray
                
                # Add text overlay
                text_shape = slide.shapes.add_textbox(
                    Inches(image_x), Inches(image_y), Inches(image_width), Inches(image_height)
                )
                text_frame = text_shape.text_frame
                text_frame.text = os.path.basename(image_path).replace('.png', '').replace('_', ' ').title()
                text_para = text_frame.paragraphs[0]
                text_para.font.name = "Segoe UI"
                text_para.font.size = Pt(12)
                text_para.font.bold = True
                text_para.font.color.rgb = RGBColor(255, 255, 255)
                text_para.alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Calculate text area
        text_start_y = y + image_height + 0.1
        text_height = height - image_height - 0.1
        text_width = width - 0.2  # Padding
        
        # Add category
        category = card_data.get('category', '')
        if category:
            # Get typography options
            category_font = card_data.get('category_font', 'Segoe UI')
            category_size = card_data.get('category_font_size', 10)
            category_color = card_data.get('category_color', '#6c757d')
            
            # Convert hex color to RGB
            category_color = category_color.lstrip('#')
            cat_r = int(category_color[0:2], 16)
            cat_g = int(category_color[2:4], 16)
            cat_b = int(category_color[4:6], 16)
            
            category_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(text_start_y), Inches(text_width), Inches(0.2)
            )
            category_frame = category_box.text_frame
            category_frame.text = category
            category_para = category_frame.paragraphs[0]
            category_para.font.name = category_font
            category_para.font.size = Pt(category_size)
            category_para.font.bold = True
            category_para.font.color.rgb = RGBColor(cat_r, cat_g, cat_b)
            category_para.alignment = PP_ALIGN.LEFT
            category_frame.word_wrap = True
        
        # Add title
        title = card_data.get('title', '')
        if title:
            # Get typography options
            title_font = card_data.get('title_font', 'Segoe UI')
            title_size = card_data.get('title_font_size', 14)
            title_color = card_data.get('title_color', '#1f4e79')
            
            # Convert hex color to RGB
            title_color = title_color.lstrip('#')
            title_r = int(title_color[0:2], 16)
            title_g = int(title_color[2:4], 16)
            title_b = int(title_color[4:6], 16)
            
            title_y = text_start_y + 0.25
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(title_y), Inches(text_width), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.name = title_font
            title_para.font.size = Pt(title_size)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(title_r, title_g, title_b)
            title_para.alignment = PP_ALIGN.LEFT
            title_frame.word_wrap = True
        
        # Add description
        description = card_data.get('description', '')
        if description:
            # Get typography options
            desc_font = card_data.get('description_font', 'Segoe UI')
            desc_size = card_data.get('description_font_size', 10)
            desc_color = card_data.get('description_color', '#6c757d')
            line_spacing = card_data.get('line_spacing', 1.2)
            
            # Convert hex color to RGB
            desc_color = desc_color.lstrip('#')
            desc_r = int(desc_color[0:2], 16)
            desc_g = int(desc_color[2:4], 16)
            desc_b = int(desc_color[4:6], 16)
            
            desc_y = title_y + 0.6
            desc_height = text_height - 0.85
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(desc_y), Inches(text_width), Inches(desc_height)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.name = desc_font
            desc_para.font.size = Pt(desc_size)
            desc_para.font.color.rgb = RGBColor(desc_r, desc_g, desc_b)
            desc_para.alignment = PP_ALIGN.LEFT
            desc_para.line_spacing = line_spacing
            desc_frame.word_wrap = True
        
        # Add badge if specified
        badge = card_data.get('badge', {})
        if badge:
            self._add_card_badge(slide, card_data, x, y, width, height, rounded_corners)
    
    def _apply_gradient_fill(self, shape, gradient_config):
        """Apply gradient fill to a shape"""
        gradient_type = gradient_config.get('type', 'linear')
        colors = gradient_config.get('colors', ['#ffffff', '#f8f9fa'])
        
        # Convert hex colors to RGB
        rgb_colors = []
        for color in colors:
            color = color.lstrip('#')
            r = int(color[0:2], 16)
            g = int(color[2:4], 16)
            b = int(color[4:6], 16)
            rgb_colors.append(RGBColor(r, g, b))
        
        # For now, use a simple two-color gradient approach
        # python-pptx has limited gradient support, so we'll use solid colors
        # and simulate gradients with color transitions
        if len(rgb_colors) >= 2:
            # Use the first color as the main fill
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb_colors[0]
            
            # Add a subtle border with the second color to simulate gradient
            shape.line.color.rgb = rgb_colors[1]
            shape.line.width = Pt(2)
        else:
            # Single color fallback
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb_colors[0]
    
    def _add_card_shadow(self, slide, x, y, width, height, rounded_corners):
        """Add shadow effect to card"""
        # Create shadow shape (slightly offset and darker)
        shadow_offset = 0.05
        shadow_x = x + shadow_offset
        shadow_y = y + shadow_offset
        
        if rounded_corners:
            shadow_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(shadow_x), Inches(shadow_y), Inches(width), Inches(height)
            )
            shadow_shape.adjustments[0] = 0.1
        else:
            shadow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(shadow_x), Inches(shadow_y), Inches(width), Inches(height)
            )
        
        # Set shadow properties
        shadow_shape.fill.solid()
        shadow_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray
        shadow_shape.line.fill.background()
        
        # Send shadow to back
        self._send_to_back(slide, shadow_shape)
    
    def _set_slide_background(self, slide, color_hex):
        """Set slide background color"""
        # Convert hex color to RGB
        color_hex = color_hex.lstrip('#')
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        
        # Create background rectangle
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.slide_width), Inches(self.slide_height)
        )
        
        # Set background color
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(r, g, b)
        bg_shape.line.fill.background()
        
        # Send to back
        self._send_to_back(slide, bg_shape)
    
    def _add_card_badge(self, slide, card_data, x, y, width, height, rounded_corners):
        """Add a badge to the card"""
        badge = card_data.get('badge', {})
        if not badge:
            return
        
        badge_text = badge.get('text', '')
        badge_color = badge.get('color', '#007bff')  # Default blue
        badge_position = badge.get('position', 'top-right')  # top-right, top-left, bottom-right, bottom-left
        badge_size = badge.get('size', 'small')  # small, medium, large
        
        # Convert hex color to RGB
        badge_color = badge_color.lstrip('#')
        badge_r = int(badge_color[0:2], 16)
        badge_g = int(badge_color[2:4], 16)
        badge_b = int(badge_color[4:6], 16)
        
        # Calculate badge dimensions based on size
        if badge_size == 'small':
            badge_width = 0.8
            badge_height = 0.3
            font_size = 8
        elif badge_size == 'medium':
            badge_width = 1.0
            badge_height = 0.4
            font_size = 10
        else:  # large
            badge_width = 1.2
            badge_height = 0.5
            font_size = 12
        
        # Calculate badge position
        if badge_position == 'top-right':
            badge_x = x + width - badge_width - 0.1
            badge_y = y + 0.1
        elif badge_position == 'top-left':
            badge_x = x + 0.1
            badge_y = y + 0.1
        elif badge_position == 'bottom-right':
            badge_x = x + width - badge_width - 0.1
            badge_y = y + height - badge_height - 0.1
        else:  # bottom-left
            badge_x = x + 0.1
            badge_y = y + height - badge_height - 0.1
        
        # Create badge background
        if rounded_corners:
            badge_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(badge_x), Inches(badge_y), Inches(badge_width), Inches(badge_height)
            )
            badge_bg.adjustments[0] = 0.3  # More rounded for badges
        else:
            badge_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(badge_x), Inches(badge_y), Inches(badge_width), Inches(badge_height)
            )
        
        badge_bg.fill.solid()
        badge_bg.fill.fore_color.rgb = RGBColor(badge_r, badge_g, badge_b)
        badge_bg.line.color.rgb = RGBColor(badge_r, badge_g, badge_b)
        badge_bg.line.width = Pt(0.5)
        
        # Add badge text
        text_shape = slide.shapes.add_textbox(
            Inches(badge_x), Inches(badge_y), Inches(badge_width), Inches(badge_height)
        )
        text_frame = text_shape.text_frame
        text_frame.text = badge_text
        text_para = text_frame.paragraphs[0]
        text_para.font.name = "Segoe UI"
        text_para.font.size = Pt(font_size)
        text_para.font.bold = True
        text_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
        text_para.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
    
    def _send_to_back(self, slide, shape):
        """Send shape to back of slide"""
        shape.element.getparent().remove(shape.element)
        slide.shapes._spTree.insert(2, shape.element)
    
    def _create_enhanced_default_slide(self, slide, slide_data, presentation_data):
        """Create enhanced default slide layout"""
        self._add_subtle_background(slide)
        self._add_enhanced_header(slide, "Enhanced Slide", presentation_data)
        self._add_enhanced_text_content(slide, {"text": "Enhanced default slide content with modern styling"})
        self._add_enhanced_footer(slide, presentation_data)


def generate_enhanced_presentations():
    """Generate enhanced presentations with visual improvements"""
    
    generator = EnhancedVisualGenerator()
    
    print("=== Enhanced PowerPoint Generation with Visual Improvements ===\n")
    
    # 1. Generate Enhanced Quarterly Business Review
    print("Generating Enhanced Quarterly Business Review...")
    qbr_dsl = BusinessTemplateExamples.create_quarterly_business_review()
    qbr_file = generator.create_presentation_from_dsl(qbr_dsl, "enhanced_quarterly_business_review.pptx")
    print(f"✅ Created: {qbr_file}")
    
    # 2. Generate Enhanced Sales Pitch
    print("Generating Enhanced Sales Pitch...")
    sales_dsl = BusinessTemplateExamples.create_sales_pitch_presentation()
    sales_file = generator.create_presentation_from_dsl(sales_dsl, "enhanced_sales_pitch.pptx")
    print(f"✅ Created: {sales_file}")
    
    # 3. Generate Enhanced Investor Pitch
    print("Generating Enhanced Investor Pitch...")
    investor_dsl = BusinessTemplateExamples.create_investor_pitch_deck()
    investor_file = generator.create_presentation_from_dsl(investor_dsl, "enhanced_investor_pitch.pptx")
    print(f"✅ Created: {investor_file}")
    
    # 4. Generate Visual Demo Presentation
    print("Generating Visual Demo Presentation...")
    demo_dsl = (BusinessDSLBuilder()
                .set_metadata(
                    title="Enhanced Template System",
                    subtitle="Visual PowerPoint Generation with 16:9 Aspect Ratio",
                    author="AI Assistant",
                    company="@cpro",
                    date=datetime.now().strftime("%B %d, %Y")
                )
                .set_theme(BusinessTheme.MODERN_MINIMAL)
                .add_title_slide()
                .add_agenda_slide("agenda", [
                    "Visual Enhancements Overview",
                    "16:9 Aspect Ratio Benefits",
                    "Hero Artwork & Graphics",
                    "Actual Charts & Data Visualization",
                    "Modern Typography & Styling",
                    "Enhanced User Experience"
                ])
                .add_section_divider("visual_section", "Visual Enhancements", "Modern Design with Professional Appeal")
                .add_content_slide(
                    "enhancements",
                    "Key Visual Improvements",
                    "bullet_list",
                    {
                        "items": [
                            {"text": "16:9 Widescreen Aspect Ratio", "level": 0},
                            {"text": "Modern standard for professional presentations", "level": 1},
                            {"text": "Better screen utilization and visual impact", "level": 1},
                            {"text": "Hero Artwork & Graphics", "level": 0},
                            {"text": "Geometric shapes and visual elements", "level": 1},
                            {"text": "Gradient backgrounds and modern styling", "level": 1},
                            {"text": "Actual Charts & Data Visualization", "level": 0},
                            {"text": "Real charts instead of placeholders", "level": 1},
                            {"text": "Professional data presentation", "level": 1},
                            {"text": "Enhanced Typography", "level": 0},
                            {"text": "Segoe UI font family for modern look", "level": 1},
                            {"text": "Improved spacing and visual hierarchy", "level": 1}
                        ]
                    }
                )
                .add_content_slide(
                    "chart_demo",
                    "Sample Chart Visualization",
                    "chart",
                    {
                        "chart_type": "bar",
                        "data": [85, 92, 78, 96, 88],
                        "labels": ["Q1", "Q2", "Q3", "Q4", "Q5"],
                        "title": "Performance Metrics Over Time"
                    }
                )
                .add_thank_you_slide(
                    contact_info={
                        "email": "enhanced@cpro.com",
                        "website": "www.cpro.com/enhanced-templates"
                    }
                )
                .build())
    
    demo_file = generator.create_presentation_from_dsl(demo_dsl, "enhanced_visual_demo.pptx")
    print(f"✅ Created: {demo_file}")
    
    return [qbr_file, sales_file, investor_file, demo_file]


if __name__ == "__main__":
    try:
        generated_files = generate_enhanced_presentations()
        
        print(f"\n🎉 Successfully generated {len(generated_files)} enhanced PowerPoint presentations!")
        print("\nEnhanced features included:")
        print("  🖼️  16:9 widescreen aspect ratio")
        print("  🎨  Hero artwork and geometric shapes")
        print("  📊  Actual charts and data visualizations")
        print("  ✨  Gradient backgrounds and modern styling")
        print("  📝  Enhanced typography with Segoe UI")
        print("  🎯  Professional visual hierarchy")
        
        print("\nGenerated files:")
        for file in generated_files:
            if os.path.exists(file):
                file_size = os.path.getsize(file) / 1024
                print(f"  📄 {file} ({file_size:.1f} KB)")
        
    except Exception as e:
        print(f"❌ Error generating enhanced presentations: {e}")
        import traceback
        traceback.print_exc()