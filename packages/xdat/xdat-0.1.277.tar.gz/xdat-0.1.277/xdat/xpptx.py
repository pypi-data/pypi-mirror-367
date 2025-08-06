"""
To modify slide layouts:
- open in WPS Presentation
- Design --> Edit slide Master
- Can right-click to rename layout
- make sure that no actual slides are created!
"""

import tempfile
from collections import Counter
import matplotlib.pyplot as plt
import pandas as pd
from slugify import slugify
import datetime as dt
from PIL import Image
import pptx
import pptx.util
from xdat import xsettings
from bidi.algorithm import get_display as fix_rtl_bidi
import re


DEFAULT_THEME = xsettings.XDAT_ROOT.joinpath('media', 'default_theme.pptx')
assert DEFAULT_THEME.exists(), DEFAULT_THEME


class Presentation:
    def __init__(self, title=None, theme=DEFAULT_THEME, print_layout=False, fake=False):
        self.theme = theme
        self.prs = pptx.Presentation(self.theme)
        self.fake = fake
        self.title = title

        if not fake:
            if print_layout:
                self.print_layout()

            if title:
                date_str = dt.datetime.now().strftime('%B %d, %Y')
                self.add_slide_title(title=title, subtitle=xsettings.PROJECT_NAME_PRETTY, note=date_str)

    def __bool__(self):
        return not self.fake

    def add_slide(self, layout_name, **kwargs):
        add_slide(self.prs, layout_name, **kwargs)

    def add_slide_title(self, title='', subtitle='', note=''):
        self.add_slide('title', title=title, subtitle=subtitle, text=note)

    def add_slide_h0(self, title=''):
        self.add_slide('main_point', title=title)

    def add_slide_h1(self, title='', subtitle=''):
        self.add_slide('section_header', title=title, subtitle=subtitle)

    def add_slide_h2(self, title='', subtitle='', desc=''):
        self.add_slide('section_title_and_description', title=title, subtitle=subtitle, text=desc)

    def add_slide_caption(self, title='', content=''):
        self.add_slide('caption', text=title, text_2=content)

    def _no_val(self, val):
        if val is None:
            return True
        if isinstance(val, str) and not val:
            return True
        return False

    def _has_val(self, val):
        return not self._no_val(val)

    def add_slide_content(self, title='', desc='', main_content='', sub_content='', sub_title=''):
        if self._no_val(desc) and self._no_val(sub_content) and self._no_val(sub_title):
            self.add_slide('title_and_body', title=title, text=main_content)
        elif self._no_val(sub_content) and self._no_val(sub_title):
            self.add_slide('left_column', title=title, text=desc, text_2=main_content)
        else:
            assert not (self._has_val(sub_content) and self._has_val(sub_title)), 'Can only specify one'
            if self._has_val(sub_title):
                self.add_slide('left_column_2', title=title, text=desc, text_2=main_content, text_3=sub_content)

            if self._has_val(sub_content):
                self.add_slide('left_column_3', title=title, text=desc, text_2=main_content, text_3=sub_content)

    def add_slide_content_2cols(self, title='', left='', right='', left_title='', right_title='', desc='', sub_content=''):
        if self._has_val(desc) or self._has_val(sub_content):
            self.add_slide('left_with_two_cols', title=title, text=desc, text_2=left, text_3=right, text_4=left_title, text_5=right_title, text_6=sub_content)
        elif self._no_val(left_title) and self._no_val(right_title):
            self.add_slide('two_columns', title=title, text=left, text_2=right)
        else:
            self.add_slide('two_columns_with_subtitles', title=title, text=left, text_2=right, text_3=left_title, text_4=right_title)

    def print_layout(self):
        print_layout(self.theme)

    def save(self, out_path=None):
        if not self.fake:
            if not out_path:
                title = self.title or 'unnamed'
                file_name = re.sub(r'[\/:*?"<>|]', ' ', title).strip()
                out_path = xsettings.OUTPUT_PATH.joinpath(f"{file_name}.pptx")

            self.prs.save(out_path)

    @classmethod
    def capture_image(cls):
        return Img()


class Img:
    """
    Everything's in inches...
    """
    DPI = 80

    def __init__(self, tight_layout=True):
        assert xsettings.CACHE_PATH is not None, "must set xsettings.CACHE_PATH"
        tmp_folder = xsettings.CACHE_PATH.joinpath('xpptx')
        tmp_folder.ensure_dir()

        self.img_path = tempfile.NamedTemporaryFile(suffix='.png', delete=False, dir=tmp_folder).name

        if tight_layout:
            plt.tight_layout()

        plt.savefig(self.img_path, pad_inches=0)
        plt.clf()
        plt.cla()
        plt.close('all')
        img = Image.open(self.img_path)
        w, h = img.size
        self.width = w/self.DPI
        self.height = h/self.DPI

    def box(self, width, height):
        rw = self.width / width
        rh = self.height / height
        rmax = max(rh, rw)
        return self.width/rmax, self.height/rmax

    def __str__(self):
        return self.img_path

    def __repr__(self):
        return self.img_path

    def __lt__(self, other):
        return False


def _slug(text):
    return slugify(text, separator="_")


def _slug_dict(input_dict):
    counts = Counter()
    out_dict = dict()
    for text, v in input_dict.items():
        try:
            parts = text.split()
            int(parts[-1])
            parts = parts[:-1]
            text = " ".join(parts)
        except:
            pass

        try:
            parts = text.split()
            if parts[-1].lower() == 'placeholder':
                parts = parts[:-1]
            text = " ".join(parts)
        except:
            pass

        text = _slug(text)
        counts[text] += 1

        if counts[text] > 1:
            text = f"{text}_{counts[text]}"

        out_dict[text] = v

    return out_dict


def _get_layout(prs):
    layouts = _slug_dict({l.name: l for l in prs.slide_layouts})
    return layouts


def _get_placeholders(slide):
    placeholders = _slug_dict({p.name: p for p in slide.placeholders})
    return placeholders


def add_slide(prs, layout_name, **kwargs):
    layouts = _get_layout(prs)
    assert layout_name in layouts, f"{layout_name} not in: {sorted(layouts)}"

    layout = layouts[layout_name]
    slide = prs.slides.add_slide(layout)

    placeholders = _get_placeholders(slide)

    for k, v in kwargs.items():
        assert k in placeholders, f"{k} not in: {sorted(placeholders)}"
        p = placeholders[k]

        if isinstance(v, dict):
            v = pd.Series(v)

        if isinstance(v, pd.Series):
            v = pd.DataFrame(v).reset_index()
            v.columns = ['', '']

        if v is None:
            p.text = ' '

        elif isinstance(v, str):
            if len(v) == 0:
                v = ' '

            p.text = v
            nlines = len(v.split('\n'))
            if nlines < 10:
                fz = None
            elif nlines < 20:
                fz = 10
            else:
                fz = 8

            if fz:
                for prg in p.text_frame.paragraphs:
                    for prgr in prg.runs:
                        prgr.font.size = pptx.util.Pt(fz)

        elif isinstance(v, list):
            if len(v) == 0:
                p.text = ' '

            else:
                tf = p.text_frame

                tf.text = v[0]
                for item in v[1:]:
                    p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0  # top‐level bullet

        elif isinstance(v, Img):
            w,h = v.box(p.width, p.height)
            w = int(w)
            h = int(h)
            slide.shapes.add_picture(str(v), p.left, p.top, height=h)
            p.text = ' '

        elif isinstance(v, pd.DataFrame):
            shape = v.shape
            p.text = ' '
            # (y, x)
            if len(v) < 10:
                fz = 10
            elif len(v) < 20:
                fz = 8
            else:
                fz = 6

            table = slide.shapes.add_table(shape[0]+1, shape[1], p.left, p.top, height=p.height, width=p.width).table
            for x in range(shape[1]):
                cell = table.cell(0, x)
                cell.text = fix_rtl_bidi(str(v.columns[x]))
                for p in cell.text_frame.paragraphs:
                    p.font.size = pptx.util.Pt(fz)

            for y in range(shape[0]):
                for x in range(shape[1]):
                    cell = table.cell(y+1, x)
                    cell.text = fix_rtl_bidi(str(v.iloc[y, x]))
                    for p in cell.text_frame.paragraphs:
                        p.font.size = pptx.util.Pt(fz)

        else:
            raise TypeError(type(v))

    for k in set(placeholders) - set(kwargs):
        p = placeholders[k]
        p.text = ' '

    return


def print_layout(template_path=None):
    prs = pptx.Presentation(template_path)
    layouts = _get_layout(prs)
    for layout_name in sorted(layouts):
        print(f"- {layout_name}")
        layout = layouts[layout_name]
        slide = prs.slides.add_slide(layout)
        placeholders = _get_placeholders(slide)
        for pl in sorted(placeholders):
            print(f"  + {pl}")

    return


if __name__ == "__main__":
    # xsettings.PROJECT_NAME = 'xdat'
    # xsettings.updated_config()
    plt.scatter([1, 2], [3, 4])
    i = Img()
    w, h = i.box(5, 5)

    prs = Presentation()

    df = pd.DataFrame({'hi': [1, 2, 3], 'there': ['a', 'b', 'c']})

    prs.add_slide('title_and_body', title='hi', text=i)

    prs.add_slide('title_and_body', title='hi', text=prs.capture_image())
    prs.add_slide('title_and_body', title='what', text=df)

    prs.save('/tmp/xdat/test.pptx')
