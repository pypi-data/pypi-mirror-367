import os

# 正则表达式（内置库）
# import re

# import comtypes.client
# from pptx import Presentation 包名为 python-pptx
# import nbformat
# import PyPDF2

# pip install comtypes -i https://pypi.tuna.tsinghua.edu.cn/simple
# pip install python-pptx -i https://pypi.tuna.tsinghua.edu.cn/simple
# pip install nbformat -i https://pypi.tuna.tsinghua.edu.cn/simple
# pip install PyPDF2 -i https://pypi.tuna.tsinghua.edu.cn/simple

def getExtension(fileName):
    """获取文件扩展名
    Args:
        fileName(str): 文件名称
    Returns:
        ext(str): 文件扩展名
    """
    return fileName[(fileName.rfind('.') + 1):].lower()

def fileParts(fileName):
    """
    拆分文件路径，获取文件路径、文件名和扩展名。

    参数:
        fileName (str): 要处理的文件路径字符串。

    返回:
        tuple: 文件路径、文件名和小写的扩展名组成的元组。
    """
    (filePath, tempFileName) = os.path.split(fileName)
    (shortName, extension) = os.path.splitext(tempFileName)
    return filePath, shortName, extension.lower()

def file_type(fileName):
    """获取文件类型

    Args:
        fileName (str): 文件名、路径、扩展名(如 .txt)

    Returns:
        str: 类型 image/office/...
    """
    _, _, ext = fileParts(fileName)

    if ext == '':
        return None

    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.tiff', '.tif', '.bmp', '.tga']:
        return 'image'

    elif ext in ['.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx']:
        return 'office'

    elif ext in ['.mp4', '.avi', '.mkv']:
        return 'video'

    elif ext in ['.mp3', '.wav', '.ogg']:
        return 'audio'

    return None



def readAllText(fileName):
    """读取文本文件全部内容，注意文件编码必须是 utf-8

    Args:
        fileName (string): 文件路径

    Returns:
        string: 文件内容
    """
    r = ''
    f = None
    try:
        f = open(fileName, 'r', encoding='utf-8')
        r = f.read()
    finally:
        if f:
            f.close()
    return r

def writeAllText(fileName, content, mode='w'):
    """
    将文本内容写入文件。

    参数:
        fileName (str): 要写入的文件路径。
        content (str): 要写入的文本内容。
        mode (str, 可选): 文件打开模式，默认为 'w'（写入模式）。

    返回:
        str: 写入文件的字符数。
    """
    r = ''
    f = None
    try:
        f = open(fileName, mode, encoding='utf-8')
        r = f.write(content)
    finally:
        if f:
            f.close()
    return r

def writeStream(fileName, content, mode='wb'):
    """
    将二进制内容写入文件。

    参数:
        fileName (str): 要写入的文件路径。
        content (bytes): 要写入的二进制内容。
        mode (str, 可选): 文件打开模式，默认为 'wb'（二进制写入模式）。

    返回:
        str: 写入文件的字符数。
    """
    r = ''
    f = None
    try:
        f = open(fileName, mode)
        r = f.write(content)
    finally:
        if f:
            f.close()
    return r

def clearAllText(fileName):
    """
    清空文件内容。

    参数:
        fileName (str): 要清空内容的文件路径。
    """
    open(fileName, 'w').close()



def readWordText(path):
    """ 获取 word 文件内容
    Args:
        path (str): 文件路径
    Returns:
        
    """
    import comtypes.client
    
    word = comtypes.client.CreateObject("Word.Application")
    word.Visible = False

    doc = word.Documents.Open(path)
    texts = []
    for paragraph in doc.Paragraphs:
        texts.append(paragraph.Range.Text)

    doc.Close(False)
    word.Quit()
    
    return ''.join(texts)

def readExcelText(path):
    """ 获取 excel 文件内容
    Args:
        path (str): 文件路径
    Returns:
        
    """
    import comtypes.client

    # 启动Excel应用
    excel = comtypes.client.CreateObject('Excel.Application')
    excel.Visible = False
    
    workbook = excel.Workbooks.Open(path)
    texts = []
    # 遍历所有工作表
    for sheet in workbook.Worksheets:
        for row in sheet.UsedRange.Rows:
            for cell in row.Cells:
                if cell.Text != '':
                    texts.append(cell.Text)
            
    workbook.Close(SaveChanges=False)
    excel.Quit()
    
    return ''.join(texts)

def readPptText(path):
    """ 获取 ppt 文件内容
    Args:
        path (str): 文件路径
    Returns:
        
    """
    from pptx import Presentation
    
    texts = []
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    
    return ''.join(texts)

def readPdfText(path):
    import PyPDF2
    reader = PyPDF2.PdfReader(path)
    texts = []
    for page in reader.pages:
        texts.append(page.extract_text())
    return ''.join(texts)



def removeEmptyLines(txt):
    # 处理正则表达式的内置库
    import re
    # 超过两个换行符改为一个
    return re.sub(r'\n+', '\n', txt)

def readNB2Txt(file_path):
    import nbformat
    # 读取.ipynb文件
    with open(file_path, 'r', encoding='utf-8') as file:
        notebook = nbformat.read(file, as_version=4)

    # 将Notebook的内容转换为纯文本字符串
    text_content = ''

    for cell in notebook.cells:
        if cell.cell_type == 'code':
            text_content += f"\n{cell.source}\n\n"
        elif cell.cell_type == 'markdown':
            text_content += f"\n{cell.source}\n\n"

    return removeEmptyLines(text_content)
    

if __name__ == '__main__':

    a = 1
    
    path, name, ext = fileParts(f'D:\测试文件\生成的文件\图片\workflow.png')
    print(path, name, ext)
    
    c = readWordText(f'D:\测试文件\office\\1.docx')
    print(c)
    c = readExcelText(f'D:\测试文件\office\\1.xlsx')
    print(c)
    c = readPptText(f'D:\测试文件\office\\1.pptx')
    print(c)
    c = readPdfText(f'D:\测试文件\office\\2.pdf')
    print(c)
