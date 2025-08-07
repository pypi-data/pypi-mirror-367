from pptx import Presentation
from .data_cleaner import clean_text  # Import the cleaner

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extracts and returns all text from a PowerPoint (.pptx) file.
    """
    prs = Presentation(file_path)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

    return clean_text("\n".join(text_runs))  # Clean the slide text